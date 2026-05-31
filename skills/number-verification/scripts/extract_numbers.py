#!/usr/bin/env python3
"""
extract_numbers.py — multi-format → JSON candidate list

Usage:
    python3 extract_numbers.py <input_file> <output_json> [--include-years]

Supports: .md .txt .docx .pptx .xlsx .pdf  (plus stdin if input_file == '-')

Output JSON schema:
{
  "source_file": "report.docx",
  "extracted_at": "2026-05-31T10:00:00",
  "total": 42,
  "candidates": [
    { "id": 1, "value": "1.2亿台", "raw": "1.2亿",
      "context": "...过去 10 年里，全球年汽车销量约 1.2 亿台...",
      "location": "page 2, paragraph 3" }
  ]
}
"""

import sys
import os
import re
import json
import argparse
from datetime import datetime

# ---------- number regex ----------
# matches numbers with optional thousand separators, decimals, sign,
# and trailing unit (中文单位 / 货币 / 百分号 / 工程单位)
NUM_PATTERN = re.compile(
    r"""
    (?<![A-Za-z0-9_\-/])                # left boundary
    (?P<sign>[-−])?                     # optional minus
    (?P<currency>[$€£¥￥])?             # currency
    (?P<num>
        \d{1,3}(?:,\d{3})+(?:\.\d+)?    # 1,234,567(.89)
        | \d+(?:\.\d+)?                 # 1234(.56)
    )
    \s*
    (?P<unit>
        %|‰|‱
        | 亿|万亿|万|千万|百万|千|百
        | 倍|个|台|辆|人|户|家|名|次|笔|张|条|平方公里|公里|公斤|吨|升|度
        | 元|美元|人民币|RMB|USD|EUR|CNY|GBP|JPY|HKD
        | bps|bp|pp|pcts?|points?|x|×|倍
        | million|billion|trillion|thousand|m|bn|tn|k
        | M|B|T|K
        | sqft|sqm|km|kg|lbs?|GW|MW|kWh|TWh|GB|TB|PB|MB|KB
    )?
    (?![A-Za-z0-9])                     # right boundary
    """,
    re.VERBOSE,
)

# fuller-width digits sometimes appear in CJK docs
FULLWIDTH_TRANS = str.maketrans("０１２３４５６７８９．，％", "0123456789.,%")

YEAR_RE = re.compile(r"^(19|20|21)\d{2}$")  # 4-digit years 1900-2199


def normalize_text(s):
    if not s:
        return ""
    return s.translate(FULLWIDTH_TRANS)


def make_value(m):
    """Reconstruct the displayed value including currency / sign / unit."""
    sign = m.group("sign") or ""
    cur = m.group("currency") or ""
    num = m.group("num")
    unit = m.group("unit") or ""
    return f"{sign}{cur}{num}{unit}".strip()


def make_raw(m):
    return m.group("num")


def is_yearlike(raw, include_years):
    if include_years:
        return False
    return bool(YEAR_RE.match(raw.replace(",", "")))


def looks_like_ordinal(text_around):
    """skip page numbers / chapter index / list ordinals"""
    ord_patterns = [
        r"第\s*\d+\s*(页|章|节|条|款|条|项|册)",
        r"page\s*\d+",
        r"chapter\s*\d+",
        r"section\s*\d+",
        r"figure\s*\d+",
        r"table\s*\d+",
        r"\bappendix\s*\d+",
        r"\bv\d+\.\d+",  # version numbers
    ]
    low = text_around.lower()
    return any(re.search(p, low) for p in ord_patterns)


def context_window(text, start, end, before=120, after=120):
    a = max(0, start - before)
    b = min(len(text), end + after)
    snippet = text[a:b].replace("\n", " ").replace("\t", " ").strip()
    snippet = re.sub(r"\s+", " ", snippet)
    return ("..." if a > 0 else "") + snippet + ("..." if b < len(text) else "")


def scan_text(text, location, include_years, sink, counter):
    text = normalize_text(text)
    if not text:
        return
    for m in NUM_PATTERN.finditer(text):
        raw = make_raw(m).replace(",", "")
        if is_yearlike(raw, include_years):
            continue
        near = text[max(0, m.start() - 24): m.end() + 24]
        if looks_like_ordinal(near):
            continue
        value = make_value(m)
        ctx = context_window(text, m.start(), m.end())
        counter[0] += 1
        sink.append({
            "id": counter[0],
            "value": value,
            "raw": raw,
            "context": ctx,
            "location": location,
        })


# ---------- per-format readers ----------

def read_text_file(path):
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def extract_md_or_txt(path, include_years):
    text = read_text_file(path)
    out = []
    counter = [0]
    # split by blank line to make location useful
    for idx, para in enumerate(re.split(r"\n\s*\n", text), 1):
        if para.strip():
            scan_text(para, f"paragraph {idx}", include_years, out, counter)
    return out


def extract_docx(path, include_years):
    try:
        from docx import Document
    except ImportError:
        sys.exit("python-docx 未安装。请 `pip install python-docx`")
    doc = Document(path)
    out = []
    counter = [0]
    for idx, para in enumerate(doc.paragraphs, 1):
        if para.text.strip():
            scan_text(para.text, f"paragraph {idx}", include_years, out, counter)
    for t_idx, table in enumerate(doc.tables, 1):
        for r_idx, row in enumerate(table.rows, 1):
            for c_idx, cell in enumerate(row.cells, 1):
                if cell.text.strip():
                    loc = f"table {t_idx} / row {r_idx} / col {c_idx}"
                    scan_text(cell.text, loc, include_years, out, counter)
    return out


def extract_pptx(path, include_years):
    try:
        from pptx import Presentation
    except ImportError:
        sys.exit("python-pptx 未安装。请 `pip install python-pptx`")
    prs = Presentation(path)
    out = []
    counter = [0]
    for s_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                txt = "\n".join(p.text for p in shape.text_frame.paragraphs)
                if txt.strip():
                    scan_text(txt, f"slide {s_idx}", include_years, out, counter)
            if getattr(shape, "has_table", False):
                tbl = shape.table
                for r_idx, row in enumerate(tbl.rows, 1):
                    for c_idx, cell in enumerate(row.cells, 1):
                        if cell.text.strip():
                            loc = f"slide {s_idx} / table row {r_idx} / col {c_idx}"
                            scan_text(cell.text, loc, include_years, out, counter)
    return out


def extract_xlsx(path, include_years):
    try:
        from openpyxl import load_workbook
    except ImportError:
        sys.exit("openpyxl 未安装。请 `pip install openpyxl`")
    wb = load_workbook(path, data_only=True, read_only=False)
    out = []
    counter = [0]
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                v = cell.value
                if v is None:
                    continue
                txt = str(v)
                if not txt.strip():
                    continue
                loc = f"{ws.title}!{cell.coordinate}"
                # numeric cell — wrap into a synthetic context so context_window sees adjacent labels
                if isinstance(v, (int, float)):
                    # find label in column A or row 1 for context
                    label_row = ws.cell(row=cell.row, column=1).value or ""
                    label_col = ws.cell(row=1, column=cell.column).value or ""
                    label_row = str(label_row).strip()
                    label_col = str(label_col).strip()
                    parts = []
                    if label_col:
                        parts.append(f"[{label_col}]")
                    if label_row:
                        parts.append(f"[{label_row}]")
                    parts.append(f"= {txt}")
                    txt = " ".join(parts)
                scan_text(txt, loc, include_years, out, counter)
    return out


def extract_pdf(path, include_years):
    try:
        import pdfplumber
    except ImportError:
        sys.exit("pdfplumber 未安装。请 `pip install pdfplumber`")
    out = []
    counter = [0]
    with pdfplumber.open(path) as pdf:
        for p_idx, page in enumerate(pdf.pages, 1):
            text = page.extract_text() or ""
            if text.strip():
                scan_text(text, f"page {p_idx}", include_years, out, counter)
    return out


def extract_stdin(include_years):
    text = sys.stdin.read()
    out = []
    counter = [0]
    scan_text(text, "inline", include_years, out, counter)
    return out


# ---------- main ----------

EXT_DISPATCH = {
    ".md":  extract_md_or_txt,
    ".markdown": extract_md_or_txt,
    ".txt": extract_md_or_txt,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".xlsx": extract_xlsx,
    ".pdf":  extract_pdf,
}


def main():
    parser = argparse.ArgumentParser(description="Extract number candidates for verification.")
    parser.add_argument("input", help="input file path (or '-' for stdin)")
    parser.add_argument("output", help="output JSON path")
    parser.add_argument("--include-years", action="store_true",
                        help="also keep 4-digit year-like numbers (1900-2199)")
    args = parser.parse_args()

    if args.input == "-":
        cands = extract_stdin(args.include_years)
        source_name = "stdin"
    else:
        path = args.input
        if not os.path.isfile(path):
            sys.exit(f"input not found: {path}")
        ext = os.path.splitext(path)[1].lower()
        if ext not in EXT_DISPATCH:
            sys.exit(f"unsupported extension: {ext}. Supported: {list(EXT_DISPATCH.keys())}")
        cands = EXT_DISPATCH[ext](path, args.include_years)
        source_name = os.path.basename(path)

    result = {
        "source_file": source_name,
        "extracted_at": datetime.now().isoformat(timespec="seconds"),
        "total": len(cands),
        "candidates": cands,
    }
    with open(args.output, "w", encoding="utf-8") as f:
        json.dump(result, f, ensure_ascii=False, indent=2)

    print(f"✓ extracted {len(cands)} number candidate(s) → {args.output}", file=sys.stderr)


if __name__ == "__main__":
    main()
